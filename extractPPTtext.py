# -*- coding: utf-8 -*-
"""
Created on Thu Oct 24 16:55:42 2019

@author: samsung
"""

# Extract all text from slides in presentation
from pptx import Presentation

path_to_presentation = 'D:\\ProgramFiles\\PythonScripts\\testPPT\\PPT1.pptx'
prs = Presentation(path_to_presentation)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)